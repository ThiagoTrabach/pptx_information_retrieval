
# coding: utf-8

# In[11]:

from pptx import Presentation
from docx import Document
from os import listdir, stat
from collections import defaultdict


# In[12]:

def get_data_pptx(document):
    
    #load the document into a python object
    prs = Presentation(document)
    
    #create output
    data = {}
    
    #retrieve metadata
    data['metadados'] = dict(author = prs.core_properties.author,
                             created = prs.core_properties.created,
                             last_modified = prs.core_properties.modified,
                             file_path = document,
                             mb_size = stat(document).st_size/1000000
                            )

    
    #retrieve text
    text_runs = []
    
    for slide in prs.slides:
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    text_runs.append(run.text)
    
    data['text'] = text_runs
    
    
    return data


# In[13]:

def get_data_docx(document):
    
    #load the document into a python object
    doc = Document(document)
    
    #create output
    data = {}
    
    #retrieve metadata
    data['metadados'] = dict(author = doc.core_properties.author,
                             created = doc.core_properties.created,
                             last_modified = doc.core_properties.modified,
                             file_path = document,
                             mb_size = stat(document).st_size/1000000
                            )
    
    
    #retrieve text
    text_runs = []
    
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            text_runs.append(run.text)
    
    data['text'] = text_runs
  
    
    return data


# Teste PPTX

# In[1]:

'''file_path = 'relatorios/' + listdir('relatorios/')[2]
dados = get_data_pptx(file_path)

dados'''


# In[2]:

'''listdir('relatorios/')[2]'''


# Teste DOCX

# In[3]:

'''
file_path = 'relatorios/' + listdir('relatorios/')[5]
dados = get_data_docx(file_path)

dados
'''


# In[ ]:



